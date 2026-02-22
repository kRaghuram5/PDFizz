"""
Utility functions for PDF conversion operations
Refactored from individual scripts to support web application
"""

import os
import shutil
import zipfile
from datetime import datetime
import PyPDF2

# Optional imports - make each dependency optional
try:
    import fitz  # PyMuPDF
    HAVE_FITZ = True
except Exception:
    fitz = None
    HAVE_FITZ = False

try:
    from pdf2docx import Converter
    HAVE_PDF2DOCX = True
except Exception:
    Converter = None
    HAVE_PDF2DOCX = False

try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    HAVE_REPORTLAB = True
except Exception:
    canvas = None
    letter = None
    HAVE_REPORTLAB = False

try:
    from PIL import Image
    HAVE_PIL = True
except Exception:
    Image = None
    HAVE_PIL = False

try:
    from docx2pdf import convert as docx2pdf_convert
    HAVE_DOCX2PDF = True
except Exception:
    docx2pdf_convert = None
    HAVE_DOCX2PDF = False

try:
    from pptx import Presentation
    HAVE_PPTX = True
except Exception:
    Presentation = None
    HAVE_PPTX = False

try:
    from openpyxl import load_workbook
    HAVE_OPENPYXL = True
except Exception:
    load_workbook = None
    HAVE_OPENPYXL = False

import subprocess
import platform


def pdf_to_word(pdf_path, output_folder, unique_id):
    """
    Convert PDF to Word document
    
    Args:
        pdf_path: Path to input PDF file
        output_folder: Directory to save output file
        unique_id: Unique identifier for the file
    
    Returns:
        Path to the generated Word file
    """
    if not HAVE_PDF2DOCX:
        raise Exception("pdf2docx is not installed. Install it with: python -m pip install pdf2docx")
    
    try:
        output_filename = f"{unique_id}_output.docx"
        output_path = os.path.join(output_folder, output_filename)
        
        converter = Converter(pdf_path)
        converter.convert(output_path, start=0, end=None)
        converter.close()
        
        return output_path
    except AttributeError as e:
        if "'Rect' object has no attribute 'get_area'" in str(e):
            raise Exception(f"PDF to Word conversion is currently unavailable due to a compatibility issue between pdf2docx (0.5.8) and PyMuPDF (1.26+). "
                          f"Alternative: Use 'PDF to Text' to extract content, then copy into Word manually. "
                          f"Technical note: pdf2docx needs an update to support newer PyMuPDF versions.")
        raise Exception(f"PDF to Word conversion failed: {str(e)}")
    except Exception as e:
        raise Exception(f"PDF to Word conversion failed: {str(e)}")


def pdf_to_text(pdf_path, output_folder, unique_id):
    """
    Extract text from PDF
    
    Args:
        pdf_path: Path to input PDF file
        output_folder: Directory to save output file
        unique_id: Unique identifier for the file
    
    Returns:
        Path to the generated text file
    """
    try:
        output_filename = f"{unique_id}_output.txt"
        output_path = os.path.join(output_folder, output_filename)
        
        with open(pdf_path, 'rb') as pdf_file:
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            with open(output_path, 'w', encoding='utf-8') as output_file:
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    output_file.write(text)
                    output_file.write('\n' + '='*80 + '\n')
        
        return output_path
    except Exception as e:
        raise Exception(f"PDF to Text conversion failed: {str(e)}")


def pdf_to_images(pdf_path, output_folder, unique_id):
    """
    Convert PDF pages to images
    
    Args:
        pdf_path: Path to input PDF file
        output_folder: Directory to save output files
        unique_id: Unique identifier for the files
    
    Returns:
        Path to the ZIP file containing all images
    """
    if not HAVE_FITZ:
        raise Exception("PyMuPDF (fitz) is not installed. Install it with: python -m pip install PyMuPDF\nOr install full requirements to enable PDF->image features.")

    try:
        # Create temporary directory for images
        temp_dir = os.path.join(output_folder, f"{unique_id}_images")
        os.makedirs(temp_dir, exist_ok=True)
        
        doc = fitz.open(pdf_path)
        for page_num, page in enumerate(doc, start=1):
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for better quality
            image_path = os.path.join(temp_dir, f"page_{page_num}.png")
            pix.save(image_path)
        
        doc.close()
        
        # Create ZIP file
        zip_filename = f"{unique_id}_images.zip"
        zip_path = os.path.join(output_folder, zip_filename)
        
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for root, dirs, files in os.walk(temp_dir):
                for file in files:
                    file_path = os.path.join(root, file)
                    zipf.write(file_path, arcname=file)
        
        # Clean up temporary directory
        shutil.rmtree(temp_dir)
        
        return zip_path
    except Exception as e:
        raise Exception(f"PDF to Images conversion failed: {str(e)}")


def word_to_pdf(word_path, output_folder, unique_id):
    """
    Convert Word document to PDF
    Uses multiple methods: docx2pdf, LibreOffice (unoconv), or pdf2docx
    
    Args:
        word_path: Path to input Word file
        output_folder: Directory to save output file
        unique_id: Unique identifier for the file
    
    Returns:
        Path to the generated PDF file
    """
    try:
        output_filename = f"{unique_id}_output.pdf"
        output_path = os.path.join(output_folder, output_filename)
        
        # Method 1: Try using LibreOffice via subprocess (most reliable in Docker)
        try:
            word_path_abs = os.path.abspath(word_path)
            output_path_abs = os.path.abspath(output_path)
            output_folder_abs = os.path.abspath(output_folder)
            
            result = subprocess.run(
                [
                    'libreoffice',
                    '--headless',
                    '--convert-to', 'pdf',
                    '--outdir', output_folder_abs,
                    word_path_abs
                ],
                capture_output=True,
                timeout=60
            )
            
            if result.returncode == 0:
                # LibreOffice creates file with original name but .pdf extension
                base_name = os.path.splitext(os.path.basename(word_path))[0]
                temp_output = os.path.join(output_folder, f"{base_name}.pdf")
                if os.path.exists(temp_output):
                    os.rename(temp_output, output_path)
                    return output_path
        except Exception as e:
            pass  # Continue to next method
        
        # Method 2: Try using docx2pdf
        if HAVE_DOCX2PDF:
            try:
                docx2pdf_convert(word_path, output_path)
                return output_path
            except Exception as e:
                pass  # Continue to next method
        
        # Method 3: Try using pdf2docx (reverse operation but works)
        if HAVE_PDF2DOCX:
            try:
                from docx import Document
                # Load the docx file and save as temp with different name
                doc = Document(word_path)
                temp_docx = os.path.join(output_folder, f"{unique_id}_temp.docx")
                doc.save(temp_docx)
                # Try converting the temp file
                docx2pdf_convert(temp_docx, output_path)
                os.remove(temp_docx)
                if os.path.exists(output_path):
                    return output_path
            except Exception as e:
                pass
        
        # Fallback to comtypes (Windows only)
        if platform.system() == "Windows":
            try:
                import comtypes.client
                word_path_abs = os.path.abspath(word_path)
                output_path_abs = os.path.abspath(output_path)
                
                word = comtypes.client.CreateObject("Word.Application")
                word.Visible = False
                doc = word.Documents.Open(word_path_abs)
                doc.SaveAs(output_path_abs, FileFormat=17)
                doc.Close()
                word.Quit()
                
                return output_path
            except Exception as e:
                pass
        
        # If all methods failed
        raise Exception("Word to PDF conversion failed. Ensure LibreOffice or docx2pdf is properly installed and configured.")
    
    except Exception as e:
        raise Exception(f"Word to PDF conversion failed: {str(e)}")


def text_to_pdf(text_path, output_folder, unique_id):
    """
    Convert text file to PDF
    
    Args:
        text_path: Path to input text file
        output_folder: Directory to save output file
        unique_id: Unique identifier for the file
    
    Returns:
        Path to the generated PDF file
    """
    if not HAVE_REPORTLAB:
        raise Exception("reportlab is not installed. Install it with: python -m pip install reportlab")
    
    try:
        output_filename = f"{unique_id}_output.pdf"
        output_path = os.path.join(output_folder, output_filename)
        
        c = canvas.Canvas(output_path, pagesize=letter)
        width, height = letter
        
        y_position = height - 50
        line_height = 15
        
        with open(text_path, 'r', encoding='utf-8') as f:
            for line in f:
                line = line.rstrip('\n')
                
                # Create new page if needed
                if y_position < 50:
                    c.showPage()
                    y_position = height - 50
                
                # Handle long lines by wrapping
                if len(line) > 80:
                    words = line.split()
                    current_line = ""
                    for word in words:
                        test_line = current_line + " " + word if current_line else word
                        if len(test_line) <= 80:
                            current_line = test_line
                        else:
                            c.drawString(50, y_position, current_line)
                            y_position -= line_height
                            current_line = word
                            if y_position < 50:
                                c.showPage()
                                y_position = height - 50
                    if current_line:
                        c.drawString(50, y_position, current_line)
                        y_position -= line_height
                else:
                    c.drawString(50, y_position, line)
                    y_position -= line_height
        
        c.save()
        return output_path
    except Exception as e:
        raise Exception(f"Text to PDF conversion failed: {str(e)}")


def images_to_pdf(image_paths, output_folder, unique_id):
    """
    Convert multiple images to a single PDF
    
    Args:
        image_paths: List of paths to input image files
        output_folder: Directory to save output file
        unique_id: Unique identifier for the file
    
    Returns:
        Path to the generated PDF file
    """
    if not HAVE_PIL:
        raise Exception("Pillow (PIL) is not installed. Install it with: python -m pip install Pillow")
    
    try:
        output_filename = f"{unique_id}_output.pdf"
        output_path = os.path.join(output_folder, output_filename)
        
        image_objects = []
        for img_path in image_paths:
            img = Image.open(img_path)
            # Convert to RGB if necessary
            if img.mode != 'RGB':
                img = img.convert('RGB')
            image_objects.append(img)
        
        if image_objects:
            image_objects[0].save(
                output_path,
                save_all=True,
                append_images=image_objects[1:] if len(image_objects) > 1 else []
            )
        
        return output_path
    except Exception as e:
        raise Exception(f"Images to PDF conversion failed: {str(e)}")


def extract_images_from_pdf(pdf_path, output_folder, unique_id):
    """
    Extract all images from a PDF file
    
    Args:
        pdf_path: Path to input PDF file
        output_folder: Directory to save output files
        unique_id: Unique identifier for the files
    
    Returns:
        Path to the ZIP file containing all extracted images
    """
    if not HAVE_FITZ:
        raise Exception("PyMuPDF (fitz) is not installed. Install it with: python -m pip install PyMuPDF\nOr install full requirements to enable image extraction features.")

    try:
        # Create temporary directory for images
        temp_dir = os.path.join(output_folder, f"{unique_id}_temp_images")
        os.makedirs(temp_dir, exist_ok=True)
        
        doc = fitz.open(pdf_path)
        image_count = 0
        
        for page_number, page in enumerate(doc, start=1):
            images = page.get_images(full=True)
            
            for img_index, img in enumerate(images, start=1):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                
                image_filename = f"page_{page_number}_image_{img_index}.{image_ext}"
                image_path = os.path.join(temp_dir, image_filename)
                
                with open(image_path, "wb") as img_file:
                    img_file.write(image_bytes)
                
                image_count += 1
        
        doc.close()
        
        if image_count == 0:
            shutil.rmtree(temp_dir)
            raise Exception("No images found in the PDF file")
        
        # Create ZIP file
        zip_filename = f"{unique_id}_extracted_images.zip"
        zip_path = os.path.join(output_folder, zip_filename)
        
        # Create ZIP with images
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for root, dirs, files in os.walk(temp_dir):
                for file in files:
                    file_path = os.path.join(root, file)
                    arcname = os.path.relpath(file_path, temp_dir)
                    zipf.write(file_path, arcname=arcname)
        
        # Clean up temporary directory
        shutil.rmtree(temp_dir)
        
        # Verify ZIP file exists
        if not os.path.exists(zip_path):
            raise Exception(f"Failed to create ZIP file at {zip_path}")
        
        print(f"Image extraction successful: {zip_path} ({image_count} images extracted)")
        return zip_path
    except Exception as e:
        print(f"Image extraction error: {str(e)}")
        raise Exception(f"Image extraction failed: {str(e)}")


def reverse_pdf(pdf_path, output_folder, unique_id):
    """
    Reverse the page order of a PDF file
    
    Args:
        pdf_path: Path to input PDF file
        output_folder: Directory to save output file
        unique_id: Unique identifier for the file
    
    Returns:
        Path to the reversed PDF file
    """
    try:
        output_filename = f"{unique_id}_reversed.pdf"
        output_path = os.path.join(output_folder, output_filename)
        
        reader = PyPDF2.PdfReader(pdf_path)
        writer = PyPDF2.PdfWriter()
        
        # Reverse order
        for page in reversed(reader.pages):
            writer.add_page(page)
        
        with open(output_path, 'wb') as output_file:
            writer.write(output_file)
        
        return output_path
    except Exception as e:
        raise Exception(f"PDF reversal failed: {str(e)}")


def merge_pdfs(pdf_paths, output_folder, unique_id):
    """
    Merge multiple PDF files into a single PDF
    
    Args:
        pdf_paths: List of paths to input PDF files
        output_folder: Directory to save output file
        unique_id: Unique identifier for the file
    
    Returns:
        Path to the merged PDF file
    """
    try:
        output_filename = f"{unique_id}_merged.pdf"
        output_path = os.path.join(output_folder, output_filename)
        
        writer = PyPDF2.PdfWriter()
        
        for pdf_path in pdf_paths:
            reader = PyPDF2.PdfReader(pdf_path)
            for page in reader.pages:
                writer.add_page(page)
        
        with open(output_path, 'wb') as output_file:
            writer.write(output_file)
        
        return output_path
    except Exception as e:
        raise Exception(f"PDF merging failed: {str(e)}")


def split_pdf(pdf_path, output_folder, unique_id, start_page, end_page):
    """
    Extract a range of pages from a PDF file
    
    Args:
        pdf_path: Path to input PDF file
        output_folder: Directory to save output file
        unique_id: Unique identifier for the file
        start_page: Starting page number (1-based)
        end_page: Ending page number (1-based, inclusive)
    
    Returns:
        Path to the extracted PDF file
    """
    try:
        output_filename = f"{unique_id}_split.pdf"
        output_path = os.path.join(output_folder, output_filename)
        
        reader = PyPDF2.PdfReader(pdf_path)
        writer = PyPDF2.PdfWriter()
        
        # Convert to 0-based indexing
        start = max(0, start_page - 1)
        end = min(len(reader.pages), end_page)
        
        for page_num in range(start, end):
            writer.add_page(reader.pages[page_num])
        
        with open(output_path, 'wb') as output_file:
            writer.write(output_file)
        
        return output_path
    except Exception as e:
        raise Exception(f"PDF splitting failed: {str(e)}")


def split_pdf_custom_ranges(pdf_path, output_folder, unique_id, ranges, merge=False):
    """
    Split a PDF into multiple PDFs based on custom page ranges.
    
    Args:
        pdf_path: Path to input PDF file
        output_folder: Directory to save output files
        unique_id: Unique identifier
        ranges: List of dicts with 'from' and 'to' keys (1-based, inclusive)
        merge: If True, merge all ranges into a single PDF instead of separate files
    
    Returns:
        Path to output file (single PDF if merge=True or single range, ZIP otherwise)
    """
    try:
        reader = PyPDF2.PdfReader(pdf_path)
        total_pages = len(reader.pages)
        
        # Create a temp directory for the split PDFs
        temp_dir = os.path.join(output_folder, f"{unique_id}_split_parts")
        os.makedirs(temp_dir, exist_ok=True)
        
        part_files = []
        for idx, r in enumerate(ranges, 1):
            start = max(0, int(r.get('from', 1)) - 1)
            end = min(total_pages, int(r.get('to', total_pages)))
            
            writer = PyPDF2.PdfWriter()
            for page_num in range(start, end):
                writer.add_page(reader.pages[page_num])
            
            part_filename = f"part_{idx}_pages_{start+1}-{end}.pdf"
            part_path = os.path.join(temp_dir, part_filename)
            with open(part_path, 'wb') as f:
                writer.write(f)
            part_files.append(part_path)
        
        # If merge flag is set, combine all ranges into one PDF
        if merge and len(part_files) > 1:
            merged_writer = PyPDF2.PdfWriter()
            for part_path in part_files:
                part_reader = PyPDF2.PdfReader(part_path)
                for page in part_reader.pages:
                    merged_writer.add_page(page)
            merged_output = os.path.join(output_folder, f"{unique_id}_split_merged.pdf")
            with open(merged_output, 'wb') as f:
                merged_writer.write(f)
            shutil.rmtree(temp_dir, ignore_errors=True)
            return merged_output
        
        # If only one range, return the single PDF
        if len(part_files) == 1:
            single_output = os.path.join(output_folder, f"{unique_id}_split.pdf")
            shutil.move(part_files[0], single_output)
            shutil.rmtree(temp_dir, ignore_errors=True)
            return single_output
        
        # Multiple ranges: zip them
        zip_path = os.path.join(output_folder, f"{unique_id}_split.zip")
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for part_path in part_files:
                zf.write(part_path, os.path.basename(part_path))
        
        shutil.rmtree(temp_dir, ignore_errors=True)
        return zip_path
    except Exception as e:
        raise Exception(f"PDF custom range split failed: {str(e)}")


def split_pdf_fixed(pdf_path, output_folder, unique_id, num_parts):
    """
    Split a PDF into N equal (or near-equal) parts.
    
    Args:
        pdf_path: Path to input PDF file
        output_folder: Directory to save output files
        unique_id: Unique identifier
        num_parts: Number of parts to split into
    
    Returns:
        Path to a ZIP file containing all split PDFs
    """
    try:
        reader = PyPDF2.PdfReader(pdf_path)
        total_pages = len(reader.pages)
        num_parts = max(1, min(int(num_parts), total_pages))
        
        pages_per_part = total_pages // num_parts
        remainder = total_pages % num_parts
        
        temp_dir = os.path.join(output_folder, f"{unique_id}_split_parts")
        os.makedirs(temp_dir, exist_ok=True)
        
        part_files = []
        current_page = 0
        for i in range(num_parts):
            # Distribute remainder pages across first parts
            part_size = pages_per_part + (1 if i < remainder else 0)
            
            writer = PyPDF2.PdfWriter()
            for j in range(part_size):
                writer.add_page(reader.pages[current_page])
                current_page += 1
            
            start_pg = current_page - part_size + 1
            end_pg = current_page
            part_filename = f"part_{i+1}_pages_{start_pg}-{end_pg}.pdf"
            part_path = os.path.join(temp_dir, part_filename)
            with open(part_path, 'wb') as f:
                writer.write(f)
            part_files.append(part_path)
        
        if len(part_files) == 1:
            single_output = os.path.join(output_folder, f"{unique_id}_split.pdf")
            shutil.move(part_files[0], single_output)
            shutil.rmtree(temp_dir, ignore_errors=True)
            return single_output
        
        zip_path = os.path.join(output_folder, f"{unique_id}_split.zip")
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for part_path in part_files:
                zf.write(part_path, os.path.basename(part_path))
        
        shutil.rmtree(temp_dir, ignore_errors=True)
        return zip_path
    except Exception as e:
        raise Exception(f"PDF fixed split failed: {str(e)}")


def split_pdf_extract_pages(pdf_path, output_folder, unique_id, pages):
    """
    Extract specific selected pages from a PDF into a new PDF.
    
    Args:
        pdf_path: Path to input PDF file
        output_folder: Directory to save output file
        unique_id: Unique identifier
        pages: List of 1-based page numbers to extract
    
    Returns:
        Path to the extracted PDF file
    """
    try:
        reader = PyPDF2.PdfReader(pdf_path)
        total_pages = len(reader.pages)
        
        writer = PyPDF2.PdfWriter()
        for page_num in sorted(set(pages)):
            if 1 <= page_num <= total_pages:
                writer.add_page(reader.pages[page_num - 1])
        
        output_filename = f"{unique_id}_split.pdf"
        output_path = os.path.join(output_folder, output_filename)
        with open(output_path, 'wb') as f:
            writer.write(f)
        
        return output_path
    except Exception as e:
        raise Exception(f"PDF page extraction failed: {str(e)}")


def compress_pdf(pdf_path, output_folder, unique_id):
    """
    Compress PDF by removing redundant streams and optimizing content
    
    Args:
        pdf_path: Path to input PDF file
        output_folder: Directory to save output file
        unique_id: Unique identifier for the file
    
    Returns:
        Path to the compressed PDF file
    """
    if not HAVE_FITZ:
        raise Exception("PyMuPDF (fitz) is not installed. Install it with: python -m pip install PyMuPDF")
    
    try:
        output_filename = f"{unique_id}_compressed.pdf"
        output_path = os.path.join(output_folder, output_filename)
        
        doc = fitz.open(pdf_path)
        
        # Compress by saving with deflate compression
        # This removes redundant content and applies stream compression
        doc.save(output_path, deflate=True, garbage=4)
        doc.close()
        
        return output_path
    except Exception as e:
        raise Exception(f"PDF compression failed: {str(e)}")


def rotate_pdf(pdf_path, output_folder, unique_id, rotation):
    """
    Rotate all pages in a PDF (supports any angle including custom angles).
    Uses PyMuPDF for arbitrary angles; falls back to PyPDF2 for multiples of 90.
    """
    import math
    try:
        angle = int(rotation) % 360
        output_filename = f"{unique_id}_rotated.pdf"
        output_path = os.path.join(output_folder, output_filename)

        if HAVE_FITZ:
            doc = fitz.open(pdf_path)
            new_doc = fitz.open()

            for page_num in range(len(doc)):
                page = doc[page_num]
                w, h = page.rect.width, page.rect.height

                # Determine new page dimensions
                if angle in (90, 270):
                    new_w, new_h = h, w
                elif angle == 180 or angle == 0:
                    new_w, new_h = w, h
                else:
                    rad = math.radians(angle)
                    new_w = abs(w * math.cos(rad)) + abs(h * math.sin(rad))
                    new_h = abs(w * math.sin(rad)) + abs(h * math.cos(rad))

                new_page = new_doc.new_page(width=new_w, height=new_h)
                new_page.show_pdf_page(new_page.rect, doc, page_num, rotate=angle)

            new_doc.save(output_path)
            new_doc.close()
            doc.close()
        else:
            # Fallback: PyPDF2 only supports multiples of 90
            snap = round(angle / 90) * 90 % 360
            reader = PyPDF2.PdfReader(pdf_path)
            writer = PyPDF2.PdfWriter()
            for page in reader.pages:
                page.rotate(snap)
                writer.add_page(page)
            with open(output_path, 'wb') as f:
                writer.write(f)

        return output_path
    except Exception as e:
        raise Exception(f"PDF rotation failed: {str(e)}")


def add_watermark(pdf_path, output_folder, unique_id, watermark_text,
                  position='center', font='helv', color='#888888', font_size=36):
    """
    Add text watermark to all pages in a PDF.
    position: top-left / top-center / top-right /
              center-left / center / center-right /
              bottom-left / bottom-center / bottom-right
    font: helv (Helvetica) | tiro (Times) | cour (Courier)
    color: hex string e.g. '#888888'
    font_size: integer
    """
    if not HAVE_FITZ:
        raise Exception("PyMuPDF (fitz) is not installed. Install it with: python -m pip install PyMuPDF")
    try:
        output_filename = f"{unique_id}_watermarked.pdf"
        output_path = os.path.join(output_folder, output_filename)

        # Parse hex color
        hex_c = color.lstrip('#')
        if len(hex_c) == 3:
            hex_c = ''.join(c*2 for c in hex_c)
        r = int(hex_c[0:2], 16) / 255
        g = int(hex_c[2:4], 16) / 255
        b = int(hex_c[4:6], 16) / 255
        rgb = (r, g, b)

        font_size = int(font_size)
        parts = position.split('-')
        vert  = parts[0]                           # top / center / bottom
        horiz = parts[1] if len(parts) > 1 else 'center'  # left / center / right

        align_map = {'left': 0, 'center': 1, 'right': 2}
        align = align_map.get(horiz, 1)

        doc = fitz.open(pdf_path)
        for page in doc:
            pw = page.rect.width
            ph = page.rect.height
            margin = 24

            if vert == 'top':
                y = margin + font_size
            elif vert == 'bottom':
                y = ph - margin
            else:  # center
                y = ph / 2

            if vert == 'center' and horiz == 'center':
                # Diagonal watermark — use TextWriter with morph for arbitrary rotation
                font_obj  = fitz.Font(fontname=font)
                tw        = fitz.TextWriter(page.rect, color=rgb)
                # Estimate half-width to roughly center the baseline
                half_w    = len(watermark_text) * font_size * 0.28
                origin    = fitz.Point(pw / 2 - half_w, ph / 2)
                tw.append(origin, watermark_text, font=font_obj, fontsize=font_size)
                pivot = fitz.Point(pw / 2, ph / 2)
                tw.write_text(page, morph=(pivot, fitz.Matrix(-45)))
            else:
                rect = fitz.Rect(margin, y - font_size - 2, pw - margin, y + 4)
                page.insert_textbox(
                    rect, watermark_text,
                    fontname=font,
                    fontsize=font_size,
                    color=rgb,
                    align=align,
                )

        doc.save(output_path)
        doc.close()
        return output_path
    except Exception as e:
        raise Exception(f"Watermark addition failed: {str(e)}")


def remove_pages(pdf_path, output_folder, unique_id, pages_to_remove):
    """
    Remove specific pages from a PDF
    
    Args:
        pdf_path: Path to input PDF file
        output_folder: Directory to save output file
        unique_id: Unique identifier for the file
        pages_to_remove: List of page numbers to remove (1-based)
    
    Returns:
        Path to the PDF file with pages removed
    """
    try:
        output_filename = f"{unique_id}_removed.pdf"
        output_path = os.path.join(output_folder, output_filename)
        
        reader = PyPDF2.PdfReader(pdf_path)
        writer = PyPDF2.PdfWriter()
        
        # Convert to set of 0-based indices
        remove_indices = set(int(p) - 1 for p in pages_to_remove)
        
        for page_num, page in enumerate(reader.pages):
            if page_num not in remove_indices:
                writer.add_page(page)
        
        with open(output_path, 'wb') as output_file:
            writer.write(output_file)
        
        return output_path
    except Exception as e:
        raise Exception(f"Removing pages failed: {str(e)}")


def pdf_to_powerpoint(pdf_path, output_folder, unique_id):
    """
    Convert PDF to PowerPoint presentation
    
    Args:
        pdf_path: Path to input PDF file
        output_folder: Directory to save output file
        unique_id: Unique identifier for the file
    
    Returns:
        Path to the output PowerPoint file
    """
    if not HAVE_PPTX:
        raise Exception("python-pptx is not installed. Install it with: pip install python-pptx")
    if not HAVE_FITZ:
        raise Exception("PyMuPDF (fitz) is not installed. Install it with: pip install PyMuPDF")
    
    try:
        output_filename = f"{unique_id}_converted.pptx"
        output_path = os.path.join(output_folder, output_filename)
        
        # Open PDF and create presentation
        pdf_doc = fitz.open(pdf_path)
        presentation = Presentation()
        
        # Set slide dimensions to match standard
        presentation.slide_width = int(10 * 914400)  # 10 inches in EMUs
        presentation.slide_height = int(7.5 * 914400)  # 7.5 inches
        
        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]
            # Convert page to image
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # High quality
            
            # Save temp image
            img_path = os.path.join(output_folder, f"temp_page_{page_num}.png")
            pix.save(img_path)
            
            # Add blank slide and insert image
            blank_layout = presentation.slide_layouts[6]  # Blank layout
            slide = presentation.slides.add_slide(blank_layout)
            slide.shapes.add_picture(img_path, 0, 0, width=presentation.slide_width, height=presentation.slide_height)
            
            # Clean up temp image
            os.remove(img_path)
        
        presentation.save(output_path)
        pdf_doc.close()
        return output_path
    except Exception as e:
        raise Exception(f"PDF to PowerPoint conversion failed: {str(e)}")


def powerpoint_to_pdf(pptx_path, output_folder, unique_id):
    """
    Convert PowerPoint to PDF using LibreOffice
    
    Args:
        pptx_path: Path to input PowerPoint file
        output_folder: Directory to save output file
        unique_id: Unique identifier for the file
    
    Returns:
        Path to the output PDF file
    """
    try:
        output_filename = f"{unique_id}_converted.pdf"
        output_path = os.path.join(output_folder, output_filename)
        
        # Convert to absolute paths
        pptx_abs = os.path.abspath(pptx_path)
        output_abs = os.path.abspath(output_folder)
        
        # Try using LibreOffice
        if platform.system() == "Windows":
            libreoffice_path = "soffice.exe"
        else:
            libreoffice_path = "libreoffice"
        
        # Convert using LibreOffice
        subprocess.run([
            libreoffice_path,
            "--headless",
            "--convert-to", "pdf",
            "--outdir", output_abs,
            pptx_abs
        ], check=True, capture_output=True)
        
        # Find the output PDF (LibreOffice creates it with original filename)
        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        source_pdf = os.path.join(output_abs, f"{base_name}.pdf")
        
        if os.path.exists(source_pdf):
            os.rename(source_pdf, output_path)
        
        return output_path
    except Exception as e:
        raise Exception(f"PowerPoint to PDF conversion failed: {str(e)}")


def excel_to_pdf(excel_path, output_folder, unique_id):
    """
    Convert Excel to PDF using LibreOffice
    
    Args:
        excel_path: Path to input Excel file
        output_folder: Directory to save output file
        unique_id: Unique identifier for the file
    
    Returns:
        Path to the output PDF file
    """
    try:
        output_filename = f"{unique_id}_converted.pdf"
        output_path = os.path.join(output_folder, output_filename)
        
        # Convert to absolute paths
        excel_abs = os.path.abspath(excel_path)
        output_abs = os.path.abspath(output_folder)
        
        # Try using LibreOffice
        if platform.system() == "Windows":
            libreoffice_path = "soffice.exe"
        else:
            libreoffice_path = "libreoffice"
        
        # Convert using LibreOffice
        subprocess.run([
            libreoffice_path,
            "--headless",
            "--convert-to", "pdf",
            "--outdir", output_abs,
            excel_abs
        ], check=True, capture_output=True)
        
        # Find the output PDF
        base_name = os.path.splitext(os.path.basename(excel_path))[0]
        source_pdf = os.path.join(output_abs, f"{base_name}.pdf")
        
        if os.path.exists(source_pdf):
            os.rename(source_pdf, output_path)
        
        return output_path
    except Exception as e:
        raise Exception(f"Excel to PDF conversion failed: {str(e)}")



def add_page_numbers(pdf_path, output_folder, unique_id, position='footer-center', font_size=10):
    """
    Add page numbers to PDF.
    position: header-left / header-center / header-right /
              footer-left / footer-center / footer-right
    font_size: integer (pt)
    """
    if not HAVE_FITZ:
        raise Exception("PyMuPDF (fitz) is not installed. Install it with: pip install PyMuPDF")
    try:
        output_filename = f"{unique_id}_numbered.pdf"
        output_path = os.path.join(output_folder, output_filename)

        font_size = int(font_size)
        parts     = position.split('-', 1)
        section   = parts[0]
        align_str = parts[1] if len(parts) > 1 else 'center'
        align_map = {'left': 0, 'center': 1, 'right': 2}
        align = align_map.get(align_str, 1)

        pdf_doc = fitz.open(pdf_path)
        margin  = 18
        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]
            pw, ph = page.rect.width, page.rect.height
            text = str(page_num + 1)

            # Compute x based on alignment
            if align_str == 'left':
                x = margin
            elif align_str == 'right':
                x = pw - margin
            else:  # center
                x = pw / 2

            # Compute y: PyMuPDF insert_text baseline is at the point y
            if section == 'header':
                y = margin + font_size
            else:  # footer
                y = ph - margin

            page.insert_text(
                fitz.Point(x, y),
                text,
                fontsize=font_size,
                color=(0, 0, 0),
            )

        pdf_doc.save(output_path)
        pdf_doc.close()
        return output_path
    except Exception as e:
        raise Exception(f"Adding page numbers failed: {str(e)}")


def repair_pdf(pdf_path, output_folder, unique_id):
    """
    Attempt to repair a damaged PDF
    
    Args:
        pdf_path: Path to input PDF file
        output_folder: Directory to save output file
        unique_id: Unique identifier for the file
    
    Returns:
        Path to the repaired PDF file
    """
    try:
        output_filename = f"{unique_id}_repaired.pdf"
        output_path = os.path.join(output_folder, output_filename)
        
        # Try PyPDF2 first for simple repair
        try:
            reader = PyPDF2.PdfReader(pdf_path)
            writer = PyPDF2.PdfWriter()
            
            # Copy all readable pages
            for page_num in range(len(reader.pages)):
                try:
                    page = reader.pages[page_num]
                    writer.add_page(page)
                except:
                    # Skip problematic pages
                    pass
            
            with open(output_path, 'wb') as output_file:
                writer.write(output_file)
            
            return output_path
        except:
            # If PyPDF2 fails, try PyMuPDF
            if not HAVE_FITZ:
                raise Exception("Could not repair PDF with available tools")
            
            pdf_doc = fitz.open(pdf_path)
            pdf_doc.save(output_path)
            pdf_doc.close()
            
            return output_path
    except Exception as e:
        raise Exception(f"PDF repair failed: {str(e)}")

    except Exception as e:
        raise Exception(f"Page removal failed: {str(e)}")
